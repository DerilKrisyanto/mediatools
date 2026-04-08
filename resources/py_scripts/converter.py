#!/usr/bin/env python3
"""
MediaTools File Converter — converter.py  v3 PRO
=================================================
Engines used per conversion type:
  PDF → DOCX  : pdf2docx  (layout-preserving, best fidelity)
  PDF → XLSX  : pdfplumber + openpyxl  (table extraction)
  PDF → PPTX  : PyMuPDF → python-pptx  (per-page image slide)
  PDF → JPG/PNG: PyMuPDF  (high-res rasterisation)
  IMG → PDF   : img2pdf   (lossless, correct page size)
  Office→ PDF : LibreOffice  (called from PHP controller directly)
  Image conv  : Pillow

Output: single JSON line on stdout
  {"success": true/false, "output": "...", "engine": "...", "error": "..."}
"""

import sys
import os
import json
import argparse
import traceback
import shutil
import tempfile

# ─────────────────────────── helpers ────────────────────────────

def ok(output: str, engine: str = "python") -> None:
    print(json.dumps({"success": True, "output": output, "engine": engine}), flush=True)
    sys.exit(0)

def fail(error: str) -> None:
    print(json.dumps({"success": False, "error": error}), flush=True)
    sys.exit(1)

def require(module_name: str, pip_name: str = None) -> object:
    import importlib
    try:
        return importlib.import_module(module_name)
    except ImportError:
        pkg = pip_name or module_name
        fail(f"Module '{module_name}' not found. Run: pip install {pkg}")


# ─────────────────────────── PDF → DOCX ─────────────────────────

def pdf_to_word(input_path: str, output_path: str) -> None:
    """
    Uses pdf2docx for layout-preserving PDF → DOCX conversion.
    Best fidelity for formatted PDFs with tables, columns, and images.
    """
    try:
        from pdf2docx import Converter
    except ImportError:
        fail("pdf2docx not installed. Run: pip install pdf2docx")

    try:
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            fail("pdf2docx produced empty output.")

        ok(output_path, "pdf2docx")
    except Exception as e:
        fail(f"pdf2docx error: {str(e)}\n{traceback.format_exc()}")


# ─────────────────────────── PDF → XLSX ─────────────────────────

def pdf_to_excel(input_path: str, output_path: str) -> None:
    """
    Extracts tables from PDF using pdfplumber, writes to XLSX.
    Each page's tables go to a separate sheet.
    """
    try:
        import pdfplumber
        import openpyxl
    except ImportError as e:
        fail(f"Missing library: {e}. Run: pip install pdfplumber openpyxl")

    try:
        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # remove default sheet

        with pdfplumber.open(input_path) as pdf:
            has_data = False
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                if tables:
                    ws = wb.create_sheet(title=f"Halaman {page_num}")
                    row_offset = 1
                    for tbl in tables:
                        for row in tbl:
                            for col_idx, cell in enumerate(row, start=1):
                                cell_val = cell if cell is not None else ""
                                ws.cell(row=row_offset, column=col_idx, value=str(cell_val))
                            row_offset += 1
                        row_offset += 1  # blank row between tables
                    has_data = True
                else:
                    # No table: extract raw text into one column
                    text = page.extract_text() or ""
                    if text.strip():
                        ws = wb.create_sheet(title=f"Halaman {page_num}")
                        for i, line in enumerate(text.split("\n"), start=1):
                            ws.cell(row=i, column=1, value=line)
                        has_data = True

            if not has_data:
                # Fallback: dump all text
                ws = wb.create_sheet(title="Konten")
                row = 1
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    for line in text.split("\n"):
                        ws.cell(row=row, column=1, value=line)
                        row += 1

        wb.save(output_path)
        ok(output_path, "pdfplumber+openpyxl")
    except Exception as e:
        fail(f"PDF→Excel error: {str(e)}\n{traceback.format_exc()}")


# ─────────────────────────── PDF → PPTX ─────────────────────────

def pdf_to_ppt(input_path: str, output_path: str) -> None:
    """
    Converts each PDF page to a high-res image, then embeds into PPTX slides.
    This preserves the exact visual layout — best approach for complex PDFs.
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from io import BytesIO
    except ImportError as e:
        fail(f"Missing library: {e}. Run: pip install PyMuPDF python-pptx")

    try:
        doc = fitz.open(input_path)
        prs = Presentation()

        # Match PPTX slide size to PDF page size (first page)
        if len(doc) > 0:
            page0    = doc[0]
            pdf_w_pt = page0.rect.width
            pdf_h_pt = page0.rect.height
            # Convert PDF points to EMU (1 pt = 12700 EMU)
            prs.slide_width  = int(pdf_w_pt * 12700)
            prs.slide_height = int(pdf_h_pt * 12700)

        blank_layout = prs.slide_layouts[6]  # blank layout

        for page_num in range(len(doc)):
            page = doc[page_num]
            # Render at 2x resolution for quality
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            img_bytes = BytesIO(pix.tobytes("png"))

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_bytes,
                left=0, top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        doc.close()
        prs.save(output_path)
        ok(output_path, "pymupdf+python-pptx")
    except Exception as e:
        fail(f"PDF→PPT error: {str(e)}\n{traceback.format_exc()}")


# ─────────────────────────── PDF → Images ───────────────────────

def pdf_to_images(input_path: str, output_dir: str, fmt: str = "jpg") -> None:
    """
    Rasterises each PDF page to JPG/PNG using PyMuPDF.
    Outputs files into output_dir as page_001.jpg, page_002.jpg, etc.
    """
    try:
        import fitz
    except ImportError:
        fail("PyMuPDF not installed. Run: pip install PyMuPDF")

    try:
        os.makedirs(output_dir, exist_ok=True)
        doc    = fitz.open(input_path)
        count  = 0
        actual = "jpeg" if fmt == "jpg" else fmt

        for i, page in enumerate(doc):
            mat = fitz.Matrix(2.5, 2.5)  # 225 DPI equivalent
            pix = page.get_pixmap(matrix=mat, alpha=False)
            fname = f"page_{i+1:03d}.{fmt}"
            fpath = os.path.join(output_dir, fname)
            pix.save(fpath)
            count += 1

        doc.close()
        if count == 0:
            fail("PDF contains no pages.")

        ok(output_dir, "pymupdf")
    except Exception as e:
        fail(f"PDF→Image error: {str(e)}\n{traceback.format_exc()}")


# ─────────────────────────── Images → PDF ───────────────────────

def images_to_pdf(input_path: str, output_path: str) -> None:
    """
    Converts JPG/PNG → PDF using img2pdf for lossless embedding.
    Falls back to Pillow if img2pdf is unavailable.
    """
    try:
        import img2pdf

        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(input_path))

        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            fail("img2pdf produced empty output.")

        ok(output_path, "img2pdf")
    except ImportError:
        # Fallback: Pillow
        try:
            from PIL import Image
            img = Image.open(input_path).convert("RGB")
            img.save(output_path, "PDF", resolution=150)
            ok(output_path, "pillow-pdf-fallback")
        except Exception as e2:
            fail(f"img2pdf not installed and Pillow fallback failed: {e2}")
    except Exception as e:
        fail(f"Image→PDF error: {str(e)}\n{traceback.format_exc()}")


# ─────────────────────────── Image conversion ───────────────────

def convert_image(input_path: str, output_path: str, target_fmt: str) -> None:
    """
    Convert between image formats (JPG ↔ PNG ↔ WebP) using Pillow.
    """
    try:
        from PIL import Image
    except ImportError:
        fail("Pillow not installed. Run: pip install Pillow")

    try:
        img = Image.open(input_path)

        # Normalise format name
        fmt_map = {
            "jpg":  ("JPEG", {"quality": 92, "optimize": True}),
            "jpeg": ("JPEG", {"quality": 92, "optimize": True}),
            "png":  ("PNG",  {"optimize": True}),
            "webp": ("WEBP", {"quality": 88, "method": 6}),
        }

        if target_fmt not in fmt_map:
            fail(f"Unsupported target format: {target_fmt}")

        pil_fmt, kwargs = fmt_map[target_fmt]

        # JPEG doesn't support alpha — convert to RGB
        if pil_fmt == "JPEG" and img.mode in ("RGBA", "P", "LA"):
            bg = Image.new("RGB", img.size, (255, 255, 255))
            if img.mode == "P":
                img = img.convert("RGBA")
            bg.paste(img, mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None)
            img = bg
        elif pil_fmt in ("PNG", "WEBP") and img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGBA")

        img.save(output_path, pil_fmt, **kwargs)
        ok(output_path, "pillow")
    except Exception as e:
        fail(f"Image conversion error: {str(e)}\n{traceback.format_exc()}")


# ─────────────────────────── MAIN ───────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="MediaTools File Converter")
    parser.add_argument("--type",      required=True, help="Conversion type")
    parser.add_argument("--input",     required=True, help="Input file path")
    parser.add_argument("--output",    required=True, help="Output file/dir path")
    parser.add_argument("--lo-binary", default="soffice", help="LibreOffice binary path")

    args = parser.parse_args()

    conv_type   = args.type
    input_path  = args.input
    output_path = args.output

    # Validate input
    if not os.path.exists(input_path):
        fail(f"Input file not found: {input_path}")

    if os.path.getsize(input_path) == 0:
        fail("Input file is empty (0 bytes).")

    # Dispatch
    if conv_type == "pdf_to_word":
        pdf_to_word(input_path, output_path)

    elif conv_type == "pdf_to_excel":
        pdf_to_excel(input_path, output_path)

    elif conv_type == "pdf_to_ppt":
        pdf_to_ppt(input_path, output_path)

    elif conv_type == "pdf_to_jpg":
        pdf_to_images(input_path, output_path, fmt="jpg")

    elif conv_type == "pdf_to_png":
        pdf_to_images(input_path, output_path, fmt="png")

    elif conv_type in ("jpg_to_pdf", "jpeg_to_pdf", "png_to_pdf"):
        images_to_pdf(input_path, output_path)

    elif conv_type == "jpg_to_png":
        convert_image(input_path, output_path, "png")

    elif conv_type == "png_to_jpg":
        convert_image(input_path, output_path, "jpg")

    elif conv_type == "jpg_to_webp":
        convert_image(input_path, output_path, "webp")

    elif conv_type == "png_to_webp":
        convert_image(input_path, output_path, "webp")

    elif conv_type == "webp_to_jpg":
        convert_image(input_path, output_path, "jpg")

    elif conv_type == "webp_to_png":
        convert_image(input_path, output_path, "png")

    else:
        fail(f"Unknown conversion type: {conv_type}")


if __name__ == "__main__":
    main()
